import os
from langchain_core.documents import Document
from docx import Document as DocxDocument  # Requires: pip install python-docx
from pptx import Presentation            # Requires: pip install python-pptx

from langchain_community.document_loaders import (
    CSVLoader, 
    UnstructuredExcelLoader, 
    BSHTMLLoader, 
    UnstructuredMarkdownLoader,
    PyMuPDFLoader
)

# === CSV Loader ===
def load_csv(path):
    """
    Loads CSV files. Creates one Document per row by default.
    """
    try:
        # csv_args can be customized for delimiters
        loader = CSVLoader(file_path=path, encoding="utf-8", csv_args={'delimiter': ','})
        return loader.load()
    except Exception as e:
        print(f"Error loading CSV {path}: {e}")
        return []

# === Excel Loader ===
def load_excel(path):
    """
    Loads .xlsx files. 
    Note: Requires 'openpyxl' installed.
    """
    try:
        # mode="elements" keeps the structure better for tables
        loader = UnstructuredExcelLoader(path, mode="elements")
        return loader.load()
    except Exception as e:
        print(f"Error loading Excel {path}: {e}")
        return []

# === HTML Loader ===
def load_html(path):
    """
    Loads HTML files using BeautifulSoup to extract just the text.
    Note: Requires 'beautifulsoup4' installed.
    """
    try:
        loader = BSHTMLLoader(path, open_encoding="utf-8")
        return loader.load()
    except Exception as e:
        print(f"Error loading HTML {path}: {e}")
        return []

# === Markdown Loader ===
def load_markdown(path):
    """
    Loads Markdown files. 
    Note: Requires 'unstructured' installed.
    """
    try:
        loader = UnstructuredMarkdownLoader(path)
        return loader.load()
    except Exception as e:
        print(f"Error loading Markdown {path}: {e}")
        return []
    
# === PDF Loader ===
def load_pdf(path):
    """Uses LangChain's PyMuPDFLoader to load PDF content."""
    try:
        loader = PyMuPDFLoader(path)
        return loader.load()
    except Exception as e:
        print(f"Error loading PDF {path}: {e}")
        return []

# === DOCX Loader ===
def load_docx(path):
    """Loads text from a .docx file using python-docx."""
    try:
        doc = DocxDocument(path)
        # Extract text from paragraphs, ignoring empty lines
        paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
        content = "\n".join(paragraphs)
        return [Document(page_content=content, metadata={"source": path})]
    except Exception as e:
        print(f"Error loading DOCX {path}: {e}")
        return []

# === PPTX Loader ===
def load_pptx(path):
    """Loads text from a .pptx file, organized by slide."""
    try:
        prs = Presentation(path)
        slides_text = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_content = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text)
            
            # Only add the slide if it contains text
            if slide_content:
                full_text = "\n".join(slide_content)
                slides_text.append(
                    Document(page_content=full_text, metadata={"source": path, "slide": slide_num})
                )
        return slides_text
    except Exception as e:
        print(f"Error loading PPTX {path}: {e}")
        return []

# === TXT Loader ===
def load_txt(path):
    """Loads text from a standard .txt file."""
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            content = f.read()
        return [Document(page_content=content, metadata={"source": path})]
    except Exception as e:
        print(f"Error loading TXT {path}: {e}")
        return []
def load_any_file(file_path):
    """Factory function to pick the correct loader based on extension."""
    ext = os.path.splitext(file_path)[1].lower()
    
    # Standard Documents
    if ext == ".pdf":
        return load_pdf(file_path)
    elif ext in [".docx", ".doc"]:
        return load_docx(file_path)
    elif ext in [".pptx", ".ppt"]:
        return load_pptx(file_path)
    elif ext == ".txt":
        return load_txt(file_path)
        
    # New Formats
    elif ext == ".csv":
        return load_csv(file_path)
    elif ext in [".xlsx", ".xls"]:
        return load_excel(file_path)
    elif ext in [".html", ".htm"]:
        return load_html(file_path)
    elif ext == ".md":
        return load_markdown(file_path)
        
    else:
        print(f"⚠️ Unsupported file type: {ext}")
        return []